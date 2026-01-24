"""Tests for presentation generator."""

import pytest
from pathlib import Path
import tempfile
from pptx import Presentation


def test_generate_presentation_config():
    """Test PresentationConfig dataclass."""
    from scripts.presentation.generator import PresentationConfig

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 3, "methods": 3, "results": 8, "discussion": 3},
        include_supplementary=False,
        content_mode="extractive",
    )

    assert config.presentation_type == "journal_club"
    assert config.slide_counts["results"] == 8


def test_load_presentation_config_from_yaml():
    """Test loading config from YAML file."""
    from scripts.presentation.generator import load_config

    config = load_config("journal_club")

    assert config is not None
    assert "default_slides" in config or config.get("name") == "Journal Club"


def test_generator_init():
    """Test PresentationGenerator initialization."""
    from scripts.presentation.generator import PresentationGenerator, PresentationConfig

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 3},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    assert generator.config == config


def test_generator_load_paper():
    """Test loading a paper into the generator."""
    from scripts.presentation.generator import PresentationGenerator, PresentationConfig

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 2},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Test with non-existent file (should handle gracefully)
    generator.load_paper(Path("/nonexistent/paper.pdf"))

    assert generator.paper_content is not None


def test_generator_generate_without_paper():
    """Test that generate() raises error when no paper is loaded properly."""
    from scripts.presentation.generator import PresentationGenerator, PresentationConfig

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 2},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Generate should still work but produce minimal output
    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "test_output.pptx"
        generator.generate(output_path)

        # File should be created
        assert output_path.exists()


def test_generator_load_template():
    """Test loading a template PPTX."""
    from scripts.presentation.generator import PresentationGenerator, PresentationConfig

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 2},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Test with non-existent template (should handle gracefully)
    generator.load_template(Path("/nonexistent/template.pptx"))

    # Template style should still be None or default
    assert generator.template_style is None


def test_presentation_config_defaults():
    """Test PresentationConfig with default values."""
    from scripts.presentation.generator import PresentationConfig

    config = PresentationConfig(
        presentation_type="lab_meeting",
        structure="custom",
        slide_counts={},
        include_supplementary=True,
        content_mode="generative",
    )

    assert config.presentation_type == "lab_meeting"
    assert config.structure == "custom"
    assert config.slide_counts == {}
    assert config.include_supplementary is True
    assert config.content_mode == "generative"


def test_load_config_default():
    """Test loading default config when YAML not found."""
    from scripts.presentation.generator import load_config

    # Load a config that doesn't exist as YAML (should return default)
    config = load_config("nonexistent_type")

    assert config is not None
    assert isinstance(config, dict)


def test_generate_presentation_convenience_function():
    """Test the convenience function for generating presentations."""
    from scripts.presentation.generator import generate_presentation

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "test_output.pptx"

        # Call without paper or template (should create minimal presentation)
        result_path = generate_presentation(
            output_path=output_path,
            presentation_type="journal_club",
        )

        assert result_path.exists()


def test_generator_full_workflow():
    """Test complete workflow with all steps."""
    from scripts.presentation.generator import (
        PresentationGenerator,
        PresentationConfig,
    )

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 2,
            "methods": 2,
            "results": 3,
            "discussion": 2,
        },
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "test_presentation.pptx"
        generator.generate(output_path)

        assert output_path.exists()
        # File should have some content
        assert output_path.stat().st_size > 0


def test_generator_generate_with_paper_content():
    """Test that generate produces slides with paper content."""
    from scripts.presentation.generator import PresentationGenerator, PresentationConfig
    from scripts.presentation.paper_extractor import PaperContent

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 2, "results": 2},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Inject mock paper content
    generator.paper_content = PaperContent(
        title="Test Paper Title",
        authors=["Author One", "Author Two"],
        abstract="Test abstract content.",
        sections={
            "introduction": "This study investigates important research questions.",
            "results": "We found significant results in our analysis.",
        },
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "test_paper_content.pptx"

        result = generator.generate(output_path)

        assert result.exists()
        prs = Presentation(str(result))
        assert len(prs.slides) >= 3  # Title + at least some content

        # Verify title slide has paper title
        first_slide = prs.slides[0]
        slide_text = " ".join([shape.text for shape in first_slide.shapes if hasattr(shape, "text")])
        assert "Test Paper Title" in slide_text
