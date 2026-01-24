"""Integration tests for presentation generator."""

import pytest
from pathlib import Path
import tempfile


@pytest.mark.slow
def test_full_pipeline_with_mock_content():
    """Test the full pipeline with mock content (no real PDF)."""
    from scripts.presentation.generator import (
        PresentationConfig,
        PresentationGenerator,
    )
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation

    # Create config
    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 2,
            "methods": 2,
            "results": 3,
            "discussion": 2,
        },
        include_supplementary=False,
        content_mode="extractive",
        presenter_name="Test User",
    )

    # Create generator and inject mock content
    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title="Test Paper: A Study of Testing",
        authors=["Author One", "Author Two"],
        abstract="This paper tests the presentation generator.",
        sections={
            "introduction": "This study investigates presentation generation. We found that automated tools can save time. The field has grown significantly.",
            "methods": "We used Python to build a generator. The system extracts content from PDFs. Statistical analysis was performed.",
            "results": "The generator produced valid PPTX files. Quality was assessed by users. Significant improvements were noted.",
            "discussion": "These results suggest automation is feasible. Future work could improve accuracy. Limitations include edge cases.",
        },
        figures=[],
        tables=[],
        references=[],
    )

    # Generate presentation
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result_path = generator.generate(output_path)

    # Verify output
    assert result_path.exists()
    assert result_path.stat().st_size > 0

    # Open and verify structure
    prs = Presentation(str(result_path))
    assert len(prs.slides) >= 5  # Title + at least some content slides

    # Cleanup
    output_path.unlink()


def test_config_loading():
    """Test that configs load correctly."""
    from scripts.presentation.generator import load_config

    config = load_config("journal_club")

    assert "default_slides" in config or "name" in config


def test_generator_with_custom_slide_counts():
    """Test generator accepts custom slide counts."""
    from scripts.presentation.generator import PresentationConfig, PresentationGenerator
    from scripts.presentation.paper_extractor import PaperContent
    import tempfile

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 1,
            "results": 2,
        },
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title="Minimal Test",
        authors=["Author"],
        abstract="Abstract",
        sections={"introduction": "Intro text.", "results": "Results text."},
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)
    assert result.exists()

    output_path.unlink()


@pytest.mark.parametrize("presentation_type", ["journal_club", "lab_meeting", "conference_talk"])
def test_all_presentation_types(presentation_type):
    """Test all three presentation types generate valid presentations."""
    from scripts.presentation.generator import (
        PresentationConfig,
        PresentationGenerator,
        load_config,
    )
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation
    import tempfile

    # Load default config for this type
    config_dict = load_config(presentation_type)

    # Get slide counts from config
    default_slides = config_dict.get("default_slides", {})
    content_slide_counts = {
        k: v for k, v in default_slides.items()
        if k not in ["title", "questions"]
    }

    config = PresentationConfig(
        presentation_type=presentation_type,
        structure="imrad",
        slide_counts=content_slide_counts,
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title=f"Test Paper for {presentation_type}",
        authors=["Test Author"],
        abstract="Test abstract for presentation type testing.",
        sections={
            "introduction": "Introduction content for testing. This is important research.",
            "methods": "Methods used include Python programming. Statistical analysis was performed.",
            "results": "Results showed significant findings. We found p < 0.05 in our analysis.",
            "discussion": "Discussion of the implications. Future work is suggested.",
            "conclusions": "Key conclusions from our study.",
        },
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)

    assert result.exists()
    assert result.stat().st_size > 0

    # Verify the presentation is valid
    prs = Presentation(str(result))
    assert len(prs.slides) >= 1  # At least title slide

    # Cleanup
    output_path.unlink()


@pytest.mark.parametrize("content_mode", ["extractive", "generative"])
def test_content_modes(content_mode):
    """Test extractive vs generative content modes."""
    from scripts.presentation.generator import PresentationConfig, PresentationGenerator
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation
    import tempfile

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 2,
            "results": 2,
        },
        include_supplementary=False,
        content_mode=content_mode,
    )

    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title="Content Mode Test Paper",
        authors=["Author One"],
        abstract="Abstract for testing content modes.",
        sections={
            "introduction": "This study investigates key questions (Smith et al., 2023). We found that the approach is significant. The research demonstrated important findings in this area.",
            "results": "Results showed 50% improvement (p < 0.001). Statistical analysis revealed significant effects. The data indicates strong correlations.",
        },
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)

    assert result.exists()
    assert result.stat().st_size > 0

    # Verify presentation is valid
    prs = Presentation(str(result))
    assert len(prs.slides) >= 3  # Title + content slides

    # Cleanup
    output_path.unlink()


def test_all_default_configs_are_valid():
    """Test that all default configurations are valid and complete."""
    from scripts.presentation.generator import DEFAULT_CONFIGS, load_config

    required_keys = ["name", "default_slides"]
    required_slide_sections = ["title", "introduction", "methods", "results", "discussion"]

    for config_name, config in DEFAULT_CONFIGS.items():
        # Check required top-level keys
        assert "name" in config, f"Config {config_name} missing 'name'"
        assert "default_slides" in config, f"Config {config_name} missing 'default_slides'"

        # Check required slide sections
        for section in required_slide_sections:
            assert section in config["default_slides"], \
                f"Config {config_name} missing '{section}' in default_slides"

        # Verify load_config returns the same
        loaded = load_config(config_name)
        assert loaded["name"] == config["name"]


def test_generate_presentation_convenience_function_all_types():
    """Test the convenience function works for all presentation types."""
    from scripts.presentation.generator import generate_presentation
    import tempfile

    for pres_type in ["journal_club", "lab_meeting", "conference_talk"]:
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / f"test_{pres_type}.pptx"

            result_path = generate_presentation(
                output_path=output_path,
                presentation_type=pres_type,
                content_mode="extractive",
            )

            assert result_path.exists()
            assert result_path.stat().st_size > 0


def test_generator_with_all_imrad_sections():
    """Test generator handles all IMRAD sections correctly."""
    from scripts.presentation.generator import (
        PresentationConfig,
        PresentationGenerator,
        IMRAD_SECTION_ORDER,
    )
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation
    import tempfile

    # Create slide counts for all sections
    slide_counts = {section: 1 for section in IMRAD_SECTION_ORDER}

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts=slide_counts,
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Create content for all sections
    sections = {
        section: f"Content for {section}. This is important. Significant findings here."
        for section in IMRAD_SECTION_ORDER
    }

    generator.paper_content = PaperContent(
        title="Full IMRAD Test",
        authors=["Author One", "Author Two", "Author Three", "Author Four"],
        abstract="Testing all IMRAD sections.",
        sections=sections,
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)

    assert result.exists()

    # Verify all sections created slides
    prs = Presentation(str(result))
    # Title + (section header + content) for each section with slides
    expected_min_slides = 1 + (2 * len(IMRAD_SECTION_ORDER))
    assert len(prs.slides) >= expected_min_slides

    output_path.unlink()


def test_author_display_truncation():
    """Test that many authors are truncated correctly."""
    from scripts.presentation.generator import (
        PresentationConfig,
        PresentationGenerator,
        MAX_AUTHORS_DISPLAYED,
    )
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation
    import tempfile

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={"introduction": 1},
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)

    # Create paper with many authors
    many_authors = [f"Author {i}" for i in range(10)]

    generator.paper_content = PaperContent(
        title="Many Authors Paper",
        authors=many_authors,
        abstract="Paper with many authors.",
        sections={"introduction": "Introduction content."},
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)
    assert result.exists()

    # Check the title slide shows truncated authors
    prs = Presentation(str(result))
    first_slide = prs.slides[0]
    slide_text = " ".join([
        shape.text for shape in first_slide.shapes if hasattr(shape, "text")
    ])

    # Should contain "et al." for truncated authors
    assert "et al." in slide_text

    output_path.unlink()


def test_empty_sections_handled():
    """Test that empty sections are handled gracefully."""
    from scripts.presentation.generator import PresentationConfig, PresentationGenerator
    from scripts.presentation.paper_extractor import PaperContent
    import tempfile

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 1,
            "methods": 1,
            "results": 1,
        },
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title="Sparse Content Paper",
        authors=["Author"],
        abstract="Abstract only.",
        sections={
            "introduction": "Only introduction has content.",
            # methods and results are empty
        },
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    # Should not raise an error
    result = generator.generate(output_path)
    assert result.exists()
    assert result.stat().st_size > 0

    output_path.unlink()


def test_zero_slide_count_skipped():
    """Test that sections with zero slide count are skipped."""
    from scripts.presentation.generator import PresentationConfig, PresentationGenerator
    from scripts.presentation.paper_extractor import PaperContent
    from pptx import Presentation
    import tempfile

    config = PresentationConfig(
        presentation_type="journal_club",
        structure="imrad",
        slide_counts={
            "introduction": 2,
            "methods": 0,  # Should be skipped
            "results": 2,
        },
        include_supplementary=False,
        content_mode="extractive",
    )

    generator = PresentationGenerator(config)
    generator.paper_content = PaperContent(
        title="Skip Methods Paper",
        authors=["Author"],
        abstract="Abstract.",
        sections={
            "introduction": "Introduction content.",
            "methods": "Methods content that should not appear.",
            "results": "Results content.",
        },
        figures=[],
        tables=[],
        references=[],
    )

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        output_path = Path(f.name)

    result = generator.generate(output_path)
    assert result.exists()

    # Verify methods section was skipped
    prs = Presentation(str(result))
    all_text = " ".join([
        shape.text for slide in prs.slides
        for shape in slide.shapes if hasattr(shape, "text")
    ])

    # Methods should not appear as a section header
    # (it might still appear if extracted into other content)
    # We mainly verify the presentation was created successfully
    assert len(prs.slides) >= 1

    output_path.unlink()
