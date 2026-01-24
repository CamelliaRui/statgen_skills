"""Tests for template parser."""

import pytest
from pathlib import Path


def test_extract_theme_colors_returns_dict():
    """Test that extract_theme_colors returns a dictionary of colors."""
    from scripts.presentation.template_parser import extract_theme_colors
    from pptx import Presentation

    # Create a minimal presentation in memory
    prs = Presentation()

    colors = extract_theme_colors(prs)

    assert isinstance(colors, dict)


def test_extract_fonts_returns_dict():
    """Test that extract_fonts returns font information."""
    from scripts.presentation.template_parser import extract_fonts
    from pptx import Presentation

    prs = Presentation()

    fonts = extract_fonts(prs)

    assert isinstance(fonts, dict)
    assert "title" in fonts
    assert "body" in fonts


def test_extract_logo_returns_none_for_empty_presentation():
    """Test that extract_logo returns None when no logo exists."""
    from scripts.presentation.template_parser import extract_logo
    from pptx import Presentation

    prs = Presentation()

    logo = extract_logo(prs)

    assert logo is None


def test_parse_template_returns_template_style():
    """Test that parse_template returns a TemplateStyle object."""
    from scripts.presentation.template_parser import parse_template, TemplateStyle
    from pptx import Presentation

    prs = Presentation()

    style = parse_template(prs)

    assert isinstance(style, TemplateStyle)
    assert hasattr(style, "colors")
    assert hasattr(style, "fonts")
    assert hasattr(style, "logo")
    assert hasattr(style, "slide_width")
    assert hasattr(style, "slide_height")