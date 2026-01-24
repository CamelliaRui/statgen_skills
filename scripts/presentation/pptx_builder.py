"""
PPTX builder for generating PowerPoint presentations.

Creates slides with content, figures, and styling based on
template and configuration.
"""

from pathlib import Path
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io

# Slide layout constants
BLANK_LAYOUT_INDEX = 6  # Index for blank slide layout in default templates


def create_presentation(
    width_inches: float = 13.333,
    height_inches: float = 7.5
) -> Presentation:
    """
    Create a new presentation with specified dimensions.

    Args:
        width_inches: Slide width (default: widescreen 16:9)
        height_inches: Slide height

    Returns:
        New Presentation object
    """
    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    return prs


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    presenter: str = "",
    date: str = "",
) -> None:
    """
    Add a title slide to the presentation.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle or paper title
        presenter: Presenter name
        date: Presentation date
    """
    slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Add presenter and date
    if presenter or date:
        info_text = ""
        if presenter:
            info_text += presenter
        if date:
            info_text += f"\n{date}" if presenter else date

        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.33), Inches(1)
        )
        info_frame = info_box.text_frame
        info_para = info_frame.paragraphs[0]
        info_para.text = info_text
        info_para.font.size = Pt(18)
        info_para.alignment = PP_ALIGN.CENTER


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    font_size: int = 18,
) -> None:
    """
    Add a content slide with bullet points.

    Args:
        prs: Presentation object
        title: Slide title
        bullets: List of bullet point strings
        font_size: Font size for bullets
    """
    slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Add bullets
    bullet_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.7)
    )
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = bullet_frame.paragraphs[0]
        else:
            para = bullet_frame.add_paragraph()

        para.text = f"• {bullet}"
        para.font.size = Pt(font_size)
        para.space_after = Pt(12)


def add_figure_slide(
    prs: Presentation,
    title: str,
    figure_bytes: bytes,
    caption: str = "",
) -> None:
    """
    Add a slide with a figure and caption.

    Args:
        prs: Presentation object
        title: Slide title
        figure_bytes: Image bytes
        caption: Figure caption
    """
    slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Add figure
    image_stream = io.BytesIO(figure_bytes)
    # Center the image
    pic = slide.shapes.add_picture(
        image_stream, Inches(1.5), Inches(1.5), width=Inches(10)
    )

    # Add caption
    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.8)
        )
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = caption
        caption_para.font.size = Pt(14)
        caption_para.font.italic = True
        caption_para.alignment = PP_ALIGN.CENTER


def add_section_header_slide(
    prs: Presentation,
    section_name: str,
) -> None:
    """
    Add a section header slide.

    Args:
        prs: Presentation object
        section_name: Name of the section
    """
    slide_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add section title centered
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(12.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = section_name
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER


def save_presentation(prs: Presentation, output_path: Path) -> None:
    """
    Save presentation to file.

    Args:
        prs: Presentation object
        output_path: Path to save the PPTX file
    """
    prs.save(str(output_path))
