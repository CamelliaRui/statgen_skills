"""
Template parser for extracting styling from PPTX files.

Extracts colors, fonts, and layout information from a template
presentation to apply to generated slides.
"""

from dataclasses import dataclass
import io
from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


@dataclass
class TemplateStyle:
    """Container for extracted template styling information."""
    colors: Dict[str, str]
    fonts: Dict[str, Dict[str, Any]]
    logo: Optional[bytes]  # PNG image bytes
    logo_position: str  # "top-left", "top-right", "bottom-left", "bottom-right"
    slide_width: float  # inches
    slide_height: float  # inches


def extract_theme_colors(prs: Presentation) -> Dict[str, str]:
    """
    Extract theme colors from a presentation.

    Args:
        prs: python-pptx Presentation object

    Returns:
        Dictionary mapping color names to hex color strings
    """
    colors = {}

    # Try to get colors from theme
    try:
        theme = prs.slide_master.theme
        if theme is not None:
            # Theme colors are typically: dk1, lt1, dk2, lt2, accent1-6
            color_scheme = theme.color_scheme
            if color_scheme is not None:
                for name in ["dk1", "lt1", "dk2", "lt2",
                            "accent1", "accent2", "accent3",
                            "accent4", "accent5", "accent6"]:
                    try:
                        color = getattr(color_scheme, name, None)
                        if color is not None:
                            colors[name] = str(color)
                    except Exception:
                        # Silently ignore: color extraction may fail for various PPTX formats
                        pass
    except Exception:
        # Silently ignore: theme access may fail for minimal or malformed presentations
        pass

    # Fallback: scan slides for commonly used colors
    if not colors:
        colors = _extract_colors_from_slides(prs)

    return colors


def _extract_colors_from_slides(prs: Presentation) -> Dict[str, str]:
    """Extract colors by scanning slide content."""
    colors = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color.rgb is not None:
                            rgb = run.font.color.rgb
                            hex_color = f"#{rgb}"
                            if "text" not in colors:
                                colors["text"] = hex_color

            # Check fill colors
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    if shape.fill.fore_color.rgb is not None:
                        rgb = shape.fill.fore_color.rgb
                        hex_color = f"#{rgb}"
                        if "fill" not in colors:
                            colors["fill"] = hex_color
                except Exception:
                    # Silently ignore: fill color extraction may fail for complex shapes
                    pass

    return colors


def extract_fonts(prs: Presentation) -> Dict[str, Dict[str, Any]]:
    """
    Extract font information from a presentation.

    Args:
        prs: python-pptx Presentation object

    Returns:
        Dictionary with 'title' and 'body' font specifications
    """
    fonts = {
        "title": {"name": "Calibri", "size": 44, "bold": True},
        "body": {"name": "Calibri", "size": 18, "bold": False},
    }

    # Try to extract from slide master
    try:
        master = prs.slide_master
        for shape in master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        run = para.runs[0]
                        font_info = {
                            "name": run.font.name or "Calibri",
                            "size": run.font.size.pt if run.font.size else 18,
                            "bold": run.font.bold or False,
                        }
                        # Heuristic: larger fonts are titles
                        if font_info["size"] and font_info["size"] > 30:
                            fonts["title"] = font_info
                        else:
                            fonts["body"] = font_info
                        break
    except Exception:
        # Silently ignore: slide master font extraction may fail; defaults are returned
        pass

    return fonts


def extract_logo(prs: Presentation) -> Optional[bytes]:
    """
    Extract logo image from presentation.

    Looks for small images in the slide master or first slide
    that appear to be logos (corner positioned, small size).

    Args:
        prs: python-pptx Presentation object

    Returns:
        PNG image bytes if logo found, None otherwise
    """
    # Check slide master first
    try:
        for shape in prs.slide_master.shapes:
            logo_bytes = _check_shape_for_logo(shape, prs)
            if logo_bytes:
                return logo_bytes
    except Exception:
        pass  # Silently ignore: slide master may not be accessible

    # Check first slide
    if prs.slides:
        for shape in prs.slides[0].shapes:
            logo_bytes = _check_shape_for_logo(shape, prs)
            if logo_bytes:
                return logo_bytes

    return None


def _check_shape_for_logo(shape, prs: Presentation) -> Optional[bytes]:
    """Check if a shape is likely a logo and extract it."""
    try:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            # Check if it's in a corner (likely logo position)
            slide_width = prs.slide_width.inches
            slide_height = prs.slide_height.inches

            left = shape.left.inches
            top = shape.top.inches
            width = shape.width.inches
            height = shape.height.inches

            # Logo heuristics: small image in corner
            is_small = width < slide_width * 0.3 and height < slide_height * 0.3
            is_corner = (left < slide_width * 0.3 or left > slide_width * 0.7) and \
                       (top < slide_height * 0.3 or top > slide_height * 0.7)

            if is_small and is_corner:
                image = shape.image
                return image.blob
    except Exception:
        pass  # Silently ignore: shape may not have image or position attributes

    return None


def parse_template(prs: Presentation) -> TemplateStyle:
    """
    Parse a template presentation and extract all styling information.

    Args:
        prs: python-pptx Presentation object

    Returns:
        TemplateStyle object with all extracted styling
    """
    colors = extract_theme_colors(prs)
    fonts = extract_fonts(prs)
    logo = extract_logo(prs)

    # Determine logo position if logo exists
    logo_position = "top-right"  # default
    if logo:
        logo_position = _detect_logo_position(prs)

    return TemplateStyle(
        colors=colors,
        fonts=fonts,
        logo=logo,
        logo_position=logo_position,
        slide_width=prs.slide_width.inches,
        slide_height=prs.slide_height.inches,
    )


def _detect_logo_position(prs: Presentation) -> str:
    """Detect where the logo is positioned."""
    try:
        for shape in prs.slide_master.shapes:
            if shape.shape_type == 13:  # Picture
                slide_width = prs.slide_width.inches
                slide_height = prs.slide_height.inches
                left = shape.left.inches
                top = shape.top.inches

                h_pos = "left" if left < slide_width / 2 else "right"
                v_pos = "top" if top < slide_height / 2 else "bottom"

                return f"{v_pos}-{h_pos}"
    except Exception:
        pass  # Silently ignore: use default position

    return "top-right"
